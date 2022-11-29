import collections 
import collections.abc
from pptx import Presentation
import aspose.slides as slides

replacements = {'Evaluation only.': '' , 'Created with Aspose.Slides for .NET Standard 2.0 22.8.':'' , 'Copyright 2004-2022Aspose Pty Ltd.':''}

def DisableLinks(filepath):
    logs = []
    print('==============Remove Links=================')
    prs = Presentation(filepath)
    slides = prs.slides
    for slide in slides:
        for shape in slide.shapes:
            log = f'[-] Remove the Link ==> {shape.click_action.hyperlink.address}'
            logs.append(log)
            print(log)
            shape.click_action.hyperlink.address = None

    prs.save(filepath)
    print('[=]Successfully remove all links.')
    print('=====================================================')
    return logs if logs != None else None 




def DisableMacro(filepath : str):
    logs = []
    dir = filepath.split('\\')
    filename = ''
    for i in range(0,len(dir) -1 ):
        filename += dir[i] +'\\'
    fn= dir[-1].split('.')
    filename +=fn+'.pptx'
    del fn,dir
    with slides.Presentation(filepath) as prs:
        if prs.vba_project.modules is None :
            return None,None
        vba_list = list(prs.vba_project.modules)
        for _ in range(len(vba_list)):
            log = f'[-] Remove the Visual Basic App ==>{prs.vba_project.modules[0].name}'
            print(log)
            logs.append(log)
            prs.vba_project.modules.remove(prs.vba_project.modules[0])
        prs.save(filename,slides.export.SaveFormat.PPTX) 
    clean_ppt(filename)
    print('[=]Successfully remove all Macros.')  
    print('==============================================================')
    return logs,filename+'.pptx' if logs != None else None,filename+'.pptx' 
    


def clean_ppt(filepath):
    global replacements
    prs = Presentation(filepath)
    for slide in prs.slides:
        for shape in slide.shapes:
            for match, replacement in replacements.items():
                if shape.has_text_frame:
                    if (shape.text.find(match)) != -1:
                        text_frame = shape.text_frame
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                cur_text = run.text
                                new_text = cur_text.replace(str(match), str(replacement))
                                run.text = new_text
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if match in cell.text:
                                new_text = cell.text.replace(match, replacement)
                                cell.text = new_text
    prs.save(filepath)


